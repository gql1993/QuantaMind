# -*- coding: utf-8 -*-
"""生成 QuantaMind 详细设计 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BG = RGBColor(0x0d, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1b, 0x22)
ACCENT = RGBColor(0x58, 0xa6, 0xff)
WHITE = RGBColor(0xe6, 0xed, 0xf3)
MUTED = RGBColor(0x8b, 0x94, 0x9e)
OK = RGBColor(0x3f, 0xb9, 0x50)
WARN = RGBColor(0xd2, 0x99, 0x22)
DANGER = RGBColor(0xf8, 0x51, 0x49)

def set_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tf

def add_bullet(tf, text, size=14, color=WHITE, level=0):
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.level = level
    return p

def add_rect(slide, left, top, width, height, fill_color=SURFACE):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color; shape.line.fill.background()
    return shape

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═══════ Slide 1: 封面 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 1, 1.2, 11, 1.2, "QuantaMind 量智大脑", 48, ACCENT, True, PP_ALIGN.CENTER)
add_text(s, 1, 2.5, 11, 0.8, "详细设计方案", 32, WHITE, False, PP_ALIGN.CENTER)
add_text(s, 1, 3.5, 11, 0.6, "量子科技自主科研 AI 中台 · 全生命周期多智能体协同平台", 18, MUTED, False, PP_ALIGN.CENTER)
add_text(s, 1, 5.0, 11, 0.8, "量子科技长三角产业创新中心\n2026年3月", 16, MUTED, False, PP_ALIGN.CENTER)

# ═══════ Slide 2: 目录 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "目录", 32, ACCENT, True)
items = ["01  发展目标与愿景", "02  发展重点", "03  发展路线图", "04  技术架构",
         "05  软件产品模块", "06  AI 智能体团队", "07  平台集成（四大域）",
         "08  工作流程", "09  数据中台与资料库", "10  关键技术指标"]
tf = add_text(s, 1, 1.2, 11, 5, items[0], 20, WHITE)
for it in items[1:]:
    add_bullet(tf, it, 20, WHITE)

# ═══════ Slide 3: 发展目标与愿景 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "01  发展目标与愿景", 32, ACCENT, True)
add_text(s, 0.5, 1.1, 12, 0.5, "打造量子芯片「设计→仿真→制造→校准→测控→数据分析」全生命周期自主科研 AI 中台", 18, WHITE)
# 三大目标
targets = [
    ("自主科研加速", "将量子芯片研发周期\n缩短 50%+", "AI Agent 自动完成\n设计→仿真→优化迭代"),
    ("全链路数据贯通", "设计/制造/测控三域\n数据实时关联追溯", "数据驱动科研决策\nAI 持续学习迭代"),
    ("多智能体协同", "12 个 AI 科学家\n覆盖全生命周期", "人机协同 + 自主发现\n7×24 不间断科研"),
]
for i, (title, desc1, desc2) in enumerate(targets):
    x = 0.3 + i * 4.3
    add_rect(s, x, 2.0, 4.0, 4.5, SURFACE)
    add_text(s, x + 0.2, 2.1, 3.6, 0.5, title, 20, ACCENT, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 2.8, 3.6, 1.5, desc1, 16, WHITE, False, PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 4.5, 3.6, 1.5, desc2, 14, MUTED, False, PP_ALIGN.CENTER)

# ═══════ Slide 4: 发展重点 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "02  发展重点", 32, ACCENT, True)
priorities = [
    ("核心能力建设", [
        "Tool Call 闭环：LLM 自主调用 107+ 工具执行操作",
        "多 Agent 协同：12 Agent 按意图路由，专属 prompt + 工具权限",
        "知识库检索：947 条知识条目，从设计文档中实时检索"]),
    ("平台深度集成", [
        "Q-EDA：Qiskit Metal + KQCircuits（已安装，真实调用）",
        "MES：CHIPMES 真实系统对接（端口 8082）",
        "测控：ARTIQ + Qiskit Pulse + Mitiq 三层测控栈",
        "数据：QCoDeS + SeaTunnel + Doris + qData 四层数据底座"]),
    ("产品化与体验", [
        "Web 客户端 12 个页面，暗色主题，OpenClaw 风格",
        "流水线引擎：3 模板（20/100/105 比特）+ 对话自动生成",
        "项目资料库：自定义文件夹 + 拖拽上传 + 自动解析",
        "全环节文件自动保存 + Word/Markdown 报告导出"]),
]
for i, (title, items) in enumerate(priorities):
    y = 1.1 + i * 2.0
    add_text(s, 0.5, y, 3, 0.4, title, 18, WARN, True)
    for j, item in enumerate(items):
        add_text(s, 3.5, y + j * 0.45, 9.5, 0.4, "· " + item, 13, WHITE if j < 2 else MUTED)

# ═══════ Slide 5: 发展路线图 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "03  发展路线图", 32, ACCENT, True)
phases = [
    ("2026 Q1-Q2\n基础平台", "#58a6ff", [
        "Gateway + Brain + Hands 核心架构",
        "Tool Call 循环 + 12 Agent",
        "Web 客户端 + 流水线引擎",
        "Q-EDA / MES 基础集成",
        "20/105 比特芯片设计验证"]),
    ("2026 Q3-Q4\nV1.0 发布", "#3fb950", [
        "真实平台全面对接",
        "Memory L1(Redis) + L4(向量RAG)",
        "知识图谱 + 文献自动追踪",
        "安全认证(JWT) + 审计日志",
        "CI/CD + Docker 容器化部署"]),
    ("2027 H1\nV2.0 深化", "#d29922", [
        "MoE 路由 + 多模型推理",
        "Heartbeat 全层级自主科研",
        "多项目 + 多团队协同",
        "训练数据集 → 模型微调",
        "千比特级芯片设计支持"]),
    ("2027 H2+\n自主科研", "#f85149", [
        "AI 自主实验设计与执行",
        "设计→制造→测控全自动",
        "黑灯工厂联动",
        "量子-经典混合优化",
        "开放生态 + 插件市场"]),
]
for i, (title, color, items) in enumerate(phases):
    x = 0.2 + i * 3.28
    c = RGBColor(int(color[1:3],16), int(color[3:5],16), int(color[5:7],16))
    add_rect(s, x, 1.2, 3.1, 5.8, SURFACE)
    add_text(s, x + 0.1, 1.3, 2.9, 0.8, title, 14, c, True, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s, x + 0.15, 2.3 + j * 0.65, 2.8, 0.6, "· " + item, 12, WHITE)

# ═══════ Slide 6: 技术架构 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "04  技术架构", 32, ACCENT, True)
layers = [
    ("表现层", "Web 客户端（12 页面）· 桌面客户端 · CLI · 项目资料库", ACCENT),
    ("网关层 Gateway", "FastAPI · 42+ REST API · WebSocket · SSE 流式 · 报告导出（Word/MD）", ACCENT),
    ("编排层 Orchestrator", "意图路由（关键词）· Tool Call 循环（≤20 轮）· 对话流水线 · 阶段审批门", ACCENT),
    ("Brain + Hands + Memory + Heartbeat", "7 LLM 提供商 | 107+ 工具 · 12 适配器 | 项目记忆 + 知识库(947条) | L0-L3 自主发现", WHITE),
    ("平台适配层 · 12 个开源项目", "Q-EDA(Metal+KQC) | MES(CHIPMES+OpenMES+secsgem) | 测控(ARTIQ+Pulse+Mitiq) | 数据(QCoDeS+SeaTunnel+Doris+qData)", WARN),
]
for i, (name, desc, color) in enumerate(layers):
    y = 1.1 + i * 1.15
    add_rect(s, 0.3, y, 12.7, 1.05, SURFACE)
    add_text(s, 0.5, y + 0.02, 3.5, 0.4, name, 15, color, True)
    add_text(s, 0.5, y + 0.45, 12.3, 0.5, desc, 11, MUTED)

# ═══════ Slide 7: 软件产品模块 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "05  软件产品模块", 32, ACCENT, True)
modules = [
    ("运转看板", "实时监控", "Gateway/任务/Agent/\n平台状态一览"),
    ("对话", "AI 交互", "SSE 流式 + Tool Call\n+ 对话流水线"),
    ("流水线", "端到端", "3 模板(20/100/105bit)\n可视化 + 审批门"),
    ("任务", "管理审批", "待办/进行/完成/审批\n模态详情"),
    ("自主发现", "Heartbeat", "L0-L3 四级心跳\n文献/实验/洞察"),
    ("Agent 团队", "12 智能体", "卡片 + 详情展开\n快捷操作 → 对话"),
    ("技能", "能力全景", "Q-EDA/MES/测控/数据\n四大平台能力"),
    ("项目资料库", "文件管理", "自定义文件夹 + 上传\n自动解析 + 持久化"),
    ("数据中台", "数据底座", "4域20表 + 血缘/质量\nAI 训练导出"),
    ("设置", "配置", "7 家 LLM 一键切换\nAPI Key 持久化"),
    ("日志", "运维", "级别筛选 + 搜索"),
    ("调试", "开发", "状态/健康 + 手动 API"),
]
for i, (name, tag, desc) in enumerate(modules):
    x = 0.2 + (i % 4) * 3.28
    y = 1.1 + (i // 4) * 2.0
    add_rect(s, x, y, 3.1, 1.8, SURFACE)
    add_text(s, x + 0.1, y + 0.1, 2.9, 0.4, name, 16, ACCENT, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, y + 0.5, 2.9, 0.3, tag, 12, WARN, False, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, y + 0.9, 2.9, 0.7, desc, 11, MUTED, False, PP_ALIGN.CENTER)

# ═══════ Slide 8: AI 智能体团队 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "06  AI 智能体团队（12 个）", 32, ACCENT, True)
agents = [
    ("协调者", "全局调度", "任务分解 · 路由 · 冲突仲裁"),
    ("芯片设计师", "Q-EDA", "Metal + KQC · GDS · QEDA 资料"),
    ("理论物理学家", "理论", "哈密顿量 · 噪声 · 方案生成"),
    ("工艺工程师", "MES 制造", "CHIPMES · 良率 · SPC · 派工"),
    ("设备运维员", "测控运维", "ARTIQ · Pulse校准 · SECS/GEM"),
    ("测控科学家", "测控科研", "T1/T2 · Mitiq纠错 · 表征"),
    ("材料科学家", "材料", "格物平台 · 高通量计算"),
    ("数据分析师", "数据", "Doris · qData · 跨域追溯"),
    ("算法工程师", "算法", "VQE/QAOA · 电路编译"),
    ("仿真工程师", "仿真", "HFSS/Q3D/Sonnet"),
    ("知识工程师", "知识", "知识库 · 文献 · 经验沉淀"),
    ("项目经理", "管理", "周报 · 里程碑 · 风险预警"),
]
for i, (name, role, desc) in enumerate(agents):
    x = 0.2 + (i % 4) * 3.28
    y = 1.1 + (i // 4) * 2.0
    add_rect(s, x, y, 3.1, 1.7, SURFACE)
    add_text(s, x + 0.1, y + 0.1, 2.9, 0.4, name, 15, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, y + 0.5, 2.9, 0.3, role, 12, ACCENT, False, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, y + 0.9, 2.9, 0.6, desc, 11, MUTED, False, PP_ALIGN.CENTER)

# ═══════ Slide 9: 平台集成 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "07  平台集成（四大域 · 13 个项目）", 32, ACCENT, True)
platforms = [
    ("Q-EDA 设计", "#58a6ff", "Qiskit Metal (IBM)\nKQCircuits (IQM)\nQEDA 团队代码", "设计 + 分析 + 版图\nLOM/EPR · GDS\n掩膜 · 16份资料"),
    ("MES 制造", "#d29922", "CHIPMES (真实系统)\nOpenMES · secsgem\nGrafana", "订单 · 工艺路线\n设备 · 良率 · SPC\nSECS/GEM 通信"),
    ("测控", "#f85149", "ARTIQ (M-Labs)\nQiskit Pulse\nMitiq (Unitary Fund)", "FPGA ns 控制\n门校准优化\nZNE/PEC/CDR 纠错"),
    ("数据中台", "#58a6ff", "QCoDeS · SeaTunnel\nDoris · qData", "采集 → ETL\nOLAP 仓库\n治理 + Text2SQL"),
]
for i, (title, color, projects, desc) in enumerate(platforms):
    x = 0.2 + i * 3.28
    c = RGBColor(int(color[1:3],16), int(color[3:5],16), int(color[5:7],16))
    add_rect(s, x, 1.1, 3.1, 5.8, SURFACE)
    add_text(s, x + 0.1, 1.2, 2.9, 0.4, title, 17, c, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 1.8, 2.9, 1.5, projects, 13, WHITE, False, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 3.8, 2.9, 1.5, desc, 12, MUTED, False, PP_ALIGN.CENTER)

# ═══════ Slide 10: 工作流程 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "08  工作流程", 32, ACCENT, True)
add_text(s, 0.5, 1.0, 12, 0.4, "以 20 比特可调耦合器芯片为例的端到端流程", 16, MUTED)
steps = [
    ("1. 芯片设计", "💎", "#58a6ff", "查阅 QEDA 资料\n查询芯片规格(真实参数)\n创建 Metal 设计\n添加 20 Xmon + 19 耦合器\n路由布线 · 导出 GDS"),
    ("2. 仿真验证", "🖥️", "#a371f7", "Q3D 电容矩阵\nEPR 本征模分析\nAnsys HFSS 导出\nSonnet 导出"),
    ("3. 工艺制造", "🏭", "#d29922", "工艺路线(8步)\n批次派工\n良率检测 · SPC\n设备告警 · 封装"),
    ("4. 校准", "🔧", "#3fb950", "ARTIQ 硬件就绪\nPulse 全套校准\n频率+振幅+DRAG+读出"),
    ("5. 测控表征", "📡", "#f85149", "光谱/Rabi/T1/T2\n读出优化\nMitiq 纠错对比\n动力学去耦"),
    ("6. 数据分析", "📊", "#58a6ff", "跨域追溯\n数据质量检查\n保存到数据中台\nAI 训练导出"),
]
for i, (title, emoji, color, desc) in enumerate(steps):
    x = 0.15 + i * 2.17
    c = RGBColor(int(color[1:3],16), int(color[3:5],16), int(color[5:7],16))
    add_rect(s, x, 1.6, 2.05, 5.2, SURFACE)
    add_text(s, x + 0.05, 1.7, 1.95, 0.5, title, 13, c, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.05, 2.3, 1.95, 4.0, desc, 11, MUTED, False, PP_ALIGN.CENTER)

# ═══════ Slide 11: 数据中台与资料库 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "09  数据中台与项目资料库", 32, ACCENT, True)
tf = add_text(s, 0.5, 1.1, 6, 5, "数据中台 · 4 域 20 表", 20, WHITE, True)
add_bullet(tf, "设计域(4表)：designs/simulations/drc/gds", 13, MUTED)
add_bullet(tf, "制造域(5表)：lots/orders/yield/spc/events", 13, MUTED)
add_bullet(tf, "测控域(5表)：characterization/calibration/benchmarks", 13, MUTED)
add_bullet(tf, "流水线域(6表)：runs/steps/params/measurements", 13, MUTED)
add_bullet(tf, "", 8)
add_bullet(tf, "跨域关联：设计→制造→测控全链路追溯", 14, WHITE)
add_bullet(tf, "数据血缘：上游来源 + 下游去向", 14, WHITE)
add_bullet(tf, "数据质量：自动检查 + 质量得分", 14, WHITE)
add_bullet(tf, "AI 训练导出：一键导出训练数据集", 14, WHITE)

tf2 = add_text(s, 7, 1.1, 6, 5, "项目资料库", 20, ACCENT, True)
add_bullet(tf2, "自定义文件夹（按团队/项目分类）", 14, WHITE)
add_bullet(tf2, "支持 11 种格式自动解析：", 14, WHITE)
add_bullet(tf2, "  Word / PDF / PPT / Excel / CSV", 12, MUTED, 1)
add_bullet(tf2, "  Python / Jupyter / GDS / JSON", 12, MUTED, 1)
add_bullet(tf2, "拖拽上传（文件 + 文件夹）", 14, WHITE)
add_bullet(tf2, "自动提取参数（GHz/nH/fF/nm…）", 14, WHITE)
add_bullet(tf2, "全环节输出文件自动保存", 14, WHITE)
add_bullet(tf2, "持久化索引（重启不丢失）", 14, WHITE)
add_bullet(tf2, "所有 Agent 均可调用查阅", 14, WHITE)

# ═══════ Slide 12: 关键技术指标 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "10  关键技术指标", 32, ACCENT, True)
metrics = [
    ("AI Agent", "12 个", "覆盖全生命周期"),
    ("注册工具", "107+", "LLM Function Calling"),
    ("API 端点", "42+", "REST + SSE + WS"),
    ("开源项目", "13 个", "四大平台集成"),
    ("LLM 提供商", "7 家", "一键热切换"),
    ("知识条目", "947 条", "21 份文档索引"),
    ("流水线模板", "3 个", "20/100/105 比特"),
    ("Web 页面", "12 个", "暗色主题 SPA"),
    ("代码规模", "4220+ 行", "34 个 Python 模块"),
    ("芯片规格", "20+105 比特", "真实设计参数"),
    ("文件格式", "11 种", "自动解析"),
    ("适配器", "13 个", "1564 行适配代码"),
]
for i, (name, value, desc) in enumerate(metrics):
    x = 0.2 + (i % 4) * 3.28
    y = 1.1 + (i // 4) * 2.0
    add_rect(s, x, y, 3.1, 1.7, SURFACE)
    add_text(s, x + 0.1, y + 0.15, 2.9, 0.4, value, 28, ACCENT, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, y + 0.7, 2.9, 0.3, name, 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.1, y + 1.1, 2.9, 0.3, desc, 12, MUTED, False, PP_ALIGN.CENTER)

# ═══════ Slide 13: 总结 ═══════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, 1, 1.0, 11, 1, "QuantaMind 量智大脑", 44, ACCENT, True, PP_ALIGN.CENTER)
add_text(s, 1, 2.2, 11, 0.6, "量子科技自主科研 AI 中台 · 详细设计方案", 22, WHITE, False, PP_ALIGN.CENTER)
summary = [
    "12 个 AI 科学家 Agent · 覆盖设计→仿真→制造→校准→测控→数据全链路",
    "107+ 工具 · LLM Function Calling · 最多 20 轮自主调用",
    "13 个开源项目集成 · Qiskit Metal / KQCircuits / CHIPMES / ARTIQ / Mitiq …",
    "3 条流水线模板（20/100/105 比特）· 可视化 + 人工审批 + Word 报告",
    "项目资料库 · 自定义文件夹 · 11 种格式自动解析 · 全环节文件保存",
    "数据中台 · 4 域 20 表 · 跨域追溯 · 知识库 947 条 · AI 训练导出",
]
tf = add_text(s, 1.5, 3.3, 10, 3.5, summary[0], 15, MUTED, False, PP_ALIGN.CENTER)
for sm in summary[1:]:
    p = add_bullet(tf, sm, 15, MUTED)
    p.alignment = PP_ALIGN.CENTER

out = "E:\\work\\QuantaMind\\docs\\QuantaMind详细设计方案.pptx"
prs.save(out)
print(f"OK: {out}")
