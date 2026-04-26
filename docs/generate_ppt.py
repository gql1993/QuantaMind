# -*- coding: utf-8 -*-
"""生成 QuantaMind 产品介绍 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

BG = RGBColor(0x0d, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1b, 0x22)
ACCENT = RGBColor(0x58, 0xa6, 0xff)
WHITE = RGBColor(0xe6, 0xed, 0xf3)
MUTED = RGBColor(0x8b, 0x94, 0x9e)
OK = RGBColor(0x3f, 0xb9, 0x50)
WARN = RGBColor(0xd2, 0x99, 0x22)
DANGER = RGBColor(0xf8, 0x51, 0x49)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_bullet(tf, text, size=14, color=WHITE, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.level = level
    return p

def add_rect(slide, left, top, width, height, fill_color=SURFACE):
    from pptx.util import Inches as I
    shape = slide.shapes.add_shape(1, I(left), I(top), I(width), I(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Slide 1: 封面 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 1.5, 11, 1.2, "QuantaMind 量智大脑", 48, ACCENT, True, PP_ALIGN.CENTER)
add_text(s, 1, 2.8, 11, 0.8, "量子科技自主科研 AI 中台", 28, WHITE, False, PP_ALIGN.CENTER)
add_text(s, 1, 4.0, 11, 0.6, "多智能体协同 · 全生命周期覆盖 · 12 个开源项目集成", 18, MUTED, False, PP_ALIGN.CENTER)
add_text(s, 1, 5.5, 11, 0.5, "量子科技长三角产业创新中心 · 2026", 14, MUTED, False, PP_ALIGN.CENTER)

# ── Slide 2: 产品定位 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "产品定位", 32, ACCENT, True)
add_text(s, 0.5, 1.1, 12, 0.5, "一个入口，覆盖量子芯片「设计 → 仿真 → 制造 → 校准 → 测控 → 数据分析」全生命周期", 18, WHITE)
# 数字卡片
cards = [("12", "AI Agent", "覆盖全链路"), ("107+", "注册工具", "可被 LLM 调用"),
         ("42+", "API 端点", "REST + SSE + WS"), ("12", "开源项目", "四大平台集成"),
         ("7", "LLM 提供商", "一键切换"), ("4", "数据域 · 20表", "全链路追溯 + 资料库")]
for i, (num, label, desc) in enumerate(cards):
    x = 0.5 + (i % 3) * 4.2
    y = 2.2 + (i // 3) * 2.2
    r = add_rect(s, x, y, 3.8, 1.8, SURFACE)
    add_text(s, x + 0.3, y + 0.2, 3.2, 0.7, num, 36, ACCENT, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.3, y + 0.9, 3.2, 0.4, label, 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.3, y + 1.3, 3.2, 0.3, desc, 12, MUTED, False, PP_ALIGN.CENTER)

# ── Slide 3: 系统架构 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "系统架构", 32, ACCENT, True)
layers = [
    ("表现层", "Web 客户端（12 页面）· 桌面客户端 · CLI", ACCENT),
    ("网关层", "FastAPI · 42 个 REST API · WebSocket · SSE 流式", ACCENT),
    ("编排层", "Orchestrator · 意图路由 · Tool Call 循环（≤20 轮）· 对话流水线", ACCENT),
    ("推理 + 工具 + 记忆 + 心跳", "Brain（7 LLM）| Hands（104 工具）| Memory（项目记忆）| Heartbeat（L0-L3）", WHITE),
    ("平台适配层", "Q-EDA（Metal+KQC）| MES（OpenMES+secsgem+Grafana）| 测控（ARTIQ+Pulse+Mitiq）| 数据（QCoDeS+SeaTunnel+Doris+qData）", WARN),
]
for i, (name, desc, color) in enumerate(layers):
    y = 1.2 + i * 1.15
    add_rect(s, 0.5, y, 12.3, 1.0, SURFACE)
    add_text(s, 0.7, y + 0.05, 3, 0.4, name, 16, color, True)
    add_text(s, 0.7, y + 0.45, 11.8, 0.4, desc, 12, MUTED)

# ── Slide 4: 12 个 Agent ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "12 个 AI 科学家 Agent", 32, ACCENT, True)
agents = [
    ("🧠", "协调者", "全局调度"), ("💎", "芯片设计师", "Q-EDA 设计"),
    ("⚛️", "理论物理学家", "理论建模"), ("🏭", "工艺工程师", "MES 制造"),
    ("🔧", "设备运维员", "测控运维"), ("📡", "测控科学家", "测控科研"),
    ("🧪", "材料科学家", "材料研发"), ("📊", "数据分析师", "数据中台"),
    ("🔬", "算法工程师", "量子算法"), ("🖥️", "仿真工程师", "电磁仿真"),
    ("📚", "知识工程师", "知识管理"), ("📋", "项目经理", "项目管理"),
]
for i, (emoji, name, role) in enumerate(agents):
    x = 0.3 + (i % 4) * 3.25
    y = 1.2 + (i // 4) * 1.9
    add_rect(s, x, y, 3.0, 1.6, SURFACE)
    add_text(s, x + 0.2, y + 0.1, 2.6, 0.5, f"{emoji} {name}", 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.2, y + 0.7, 2.6, 0.3, role, 13, ACCENT, False, PP_ALIGN.CENTER)

# ── Slide 5: 四大平台集成 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "四大平台 · 12 个开源项目集成", 32, ACCENT, True)
platforms = [
    ("Q-EDA 设计", "Qiskit Metal（IBM）\nKQCircuits（IQM）", "设计+分析+版图+掩膜\n✅ 已安装", OK),
    ("MES 制造", "OpenMES + secsgem\n+ Grafana", "工艺/批次/良率/SPC\nSECS/GEM 设备通信", WARN),
    ("测控", "ARTIQ + Qiskit Pulse\n+ Mitiq", "实时控制 → 校准\n→ 错误缓解", DANGER),
    ("数据中台", "QCoDeS + SeaTunnel\n+ Doris + qData", "采集 → ETL → OLAP\n→ 治理", ACCENT),
]
for i, (title, projects, desc, color) in enumerate(platforms):
    x = 0.3 + i * 3.25
    add_rect(s, x, 1.2, 3.0, 5.5, SURFACE)
    add_text(s, x + 0.15, 1.3, 2.7, 0.5, title, 18, color, True, PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 2.0, 2.7, 1.0, projects, 13, WHITE, False, PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 3.5, 2.7, 1.0, desc, 12, MUTED, False, PP_ALIGN.CENTER)

# ── Slide 6: Tool Call 循环 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "核心机制：Tool Call 循环", 32, ACCENT, True)
add_text(s, 0.5, 1.1, 12, 0.5, "LLM 通过 Function Calling 自主调用 104 个工具，最多 20 轮", 18, WHITE)
steps_text = [
    "1. 用户发送消息 → Orchestrator 按关键词路由到最优 Agent",
    "2. 加载 Agent 专属 system prompt + 允许的工具列表",
    "3. 注入项目记忆（Memory）→ 发给 LLM（带 tools 参数）",
    "4. LLM 返回 tool_calls → Hands 执行工具 → 结果回注对话",
    "5. LLM 根据结果继续（可能再次调用工具）→ 循环最多 20 轮",
    "6. 最终回复流式返回给客户端 + 自动生成对话流水线",
]
tf = add_text(s, 0.8, 2.0, 11, 4, steps_text[0], 16, WHITE)
for st in steps_text[1:]:
    add_bullet(tf, st, 16, WHITE)

# ── Slide 7: 流水线引擎 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "流水线引擎 · 端到端可视化", 32, ACCENT, True)
tf = add_text(s, 0.5, 1.1, 6, 5, "两种流水线", 20, WHITE, True)
add_bullet(tf, "模板流水线（PL-）：预定义端到端流程", 15, WHITE)
add_bullet(tf, "  如 20 比特芯片：6 阶段 101 步 7 Agent", 13, MUTED, 1)
add_bullet(tf, "对话流水线（CL-）：对话中自动生成", 15, WHITE)
add_bullet(tf, "  每次工具调用自动记录并可视化", 13, MUTED, 1)
add_bullet(tf, "", 8, WHITE)
add_bullet(tf, "人工干预", 20, WHITE)
add_bullet(tf, "暂停 / 恢复 / 终止 / 阶段审批门", 15, WHITE)
add_bullet(tf, "", 8, WHITE)
add_bullet(tf, "报告导出", 20, WHITE)
add_bullet(tf, "Word + Markdown：含芯片规格 + 每步数据 + 中文解释", 15, WHITE)

tf2 = add_text(s, 7, 1.1, 6, 5, "20 比特流水线阶段", 20, ACCENT, True)
stages = ["1. 💎 芯片设计（64 步）", "2. 🖥️ 仿真验证（8 步）", "3. 🏭 工艺制造（9 步）",
          "4. 🔧 设备校准（4 步）", "5. 📡 测控表征（9 步）", "6. 📊 数据分析（7 步）"]
for st in stages:
    add_bullet(tf2, st, 16, WHITE)

# ── Slide 8: 数据中台 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "数据中台 · 全链路追溯", 32, ACCENT, True)
add_text(s, 0.5, 1.1, 12, 0.5, "Q-EDA(设计) + MES(制造) + QCoDeS(测控) → SeaTunnel → Doris → qData → AI Agent", 16, MUTED)
domains = [
    ("设计域", "4 表", "chip_designs / simulation_results / drc_reports / gds_exports"),
    ("制造域", "5 表", "lots / work_orders / yield_records / spc_data / equipment_events"),
    ("测控域", "5 表", "qubit_characterization / calibration_records / gate_benchmarks / error_mitigation / raw_measurements"),
    ("流水线域", "6 表", "pipeline_runs / pipeline_steps / design_params / calibration_params / measurement_results / fabrication_records"),
]
for i, (name, count, tables) in enumerate(domains):
    y = 2.0 + i * 1.25
    add_rect(s, 0.5, y, 12.3, 1.1, SURFACE)
    add_text(s, 0.7, y + 0.05, 2, 0.4, f"{name}（{count}）", 16, ACCENT, True)
    add_text(s, 0.7, y + 0.5, 11.8, 0.4, tables, 12, MUTED)

add_text(s, 0.5, 6.6, 12, 0.4, "每步数据实时保存 · AI 训练数据集一键导出 · 跨域关联查询 · 数据血缘追溯 · 数据质量检查", 14, OK)

# ── Slide 9: 20 比特芯片案例 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "实战案例：20 比特可调耦合器量子芯片", 28, ACCENT, True)
specs = [
    "文档编号：TGQ-200-000-FA09-2025",
    "芯片尺寸：12.5 × 12.5 mm · 蓝宝石衬底",
    "20 个 Xmon 量子比特（Q_odd 4.65GHz / Q_even 4.56GHz）",
    "19 个可调耦合器（6.88GHz · Transmon + SQUID）",
    "约瑟夫森结：铝-氧化铝-铝 · 200×200nm · Dolan Bridge",
    "设计目标：单门 ≥99.9% · 双门 ≥99% · 读出 ≥99% · T1≥20μs",
    "端到端流水线：6 阶段 · 7 Agent · 101 步 · ~6 秒完成",
]
tf = add_text(s, 0.8, 1.2, 11, 5, specs[0], 16, WHITE)
for sp in specs[1:]:
    add_bullet(tf, sp, 16, WHITE)

# ── Slide 10: Web 客户端 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "Web 客户端 · 12 个页面", 32, ACCENT, True)
pages = [
    ("运转看板", "统计卡片可点击跳转，实时刷新"),
    ("对话", "SSE 流式 + 对话流水线面板"),
    ("流水线", "模板列表 + 历史 + 实时可视化"),
    ("任务", "Tab 筛选 + 模态详情 + 审批"),
    ("自主发现", "Heartbeat 产出，可跳转对话"),
    ("Agent 团队", "12 Agent 卡片 + 详情展开"),
    ("技能", "四大平台能力全景 + 触发按钮"),
    ("数据中台", "域表结构 + 血缘 + 质量 + AI 导出"),
    ("会话管理", "按首条消息命名 + 删除"),
    ("日志", "级别筛选 + 文本搜索"),
    ("调试", "系统状态 + 手动 API 调用"),
    ("设置", "7 家 LLM 一键切换 + 持久化"),
]
for i, (name, desc) in enumerate(pages):
    x = 0.3 + (i % 3) * 4.3
    y = 1.2 + (i // 3) * 1.4
    add_rect(s, x, y, 4.0, 1.2, SURFACE)
    add_text(s, x + 0.15, y + 0.1, 3.7, 0.4, name, 16, ACCENT, True)
    add_text(s, x + 0.15, y + 0.55, 3.7, 0.4, desc, 12, MUTED)

# ── Slide 11: 技术亮点 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0.5, 0.3, 12, 0.7, "技术亮点", 32, ACCENT, True)
highlights = [
    ("Function Calling 闭环", "LLM 自主调用 104 工具 → 执行 → 结果回注 → 继续推理，最多 20 轮"),
    ("多 Agent 协同", "12 Agent 按关键词路由，每个 Agent 有专属 prompt 和工具权限"),
    ("优雅降级", "12 个开源库未安装时自动 Mock，不影响系统运行"),
    ("实时流水线", "对话自动生成可视化流水线，每步实时保存到数据中台"),
    ("LLM 热切换", "Web 页面一键切换 7 家 LLM 提供商，配置持久化"),
    ("全链路数据", "设计→制造→测控三域数据在 Doris 中跨域关联追溯"),
]
for i, (title, desc) in enumerate(highlights):
    y = 1.2 + i * 0.95
    add_text(s, 0.8, y, 3, 0.4, "✦ " + title, 18, OK, True)
    add_text(s, 4, y, 9, 0.4, desc, 14, MUTED)

# ── Slide 12: 总结 ──
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 1.5, 11, 1, "QuantaMind 量智大脑", 44, ACCENT, True, PP_ALIGN.CENTER)
add_text(s, 1, 2.8, 11, 0.6, "量子科技自主科研 AI 中台", 24, WHITE, False, PP_ALIGN.CENTER)
summary = [
    "12 个 AI 科学家 Agent · 覆盖量子芯片全生命周期",
    "104 个工具 · LLM Function Calling 自主调用",
    "12 个开源项目集成 · Qiskit Metal / KQCircuits / ARTIQ / Mitiq …",
    "端到端流水线 · 可视化 + 人工审批 + Word 报告",
    "数据中台 · 4 域 20 表 · 跨域追溯 · AI 训练数据导出",
]
tf = add_text(s, 1.5, 3.8, 10, 3, summary[0], 16, MUTED, False, PP_ALIGN.CENTER)
for sm in summary[1:]:
    p = add_bullet(tf, sm, 16, MUTED)
    p.alignment = PP_ALIGN.CENTER

out = "E:\\work\\QuantaMind\\docs\\QuantaMind产品介绍.pptx"
prs.save(out)
print(f"OK: {out}")
