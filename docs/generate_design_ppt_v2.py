# -*- coding: utf-8 -*-
"""生成 QuantaMind 详细设计 PPT V2 — 含架构框图"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

BG = RGBColor(0xFF, 0xFF, 0xFF)  # 白色背景（正式风格）
TITLE_BG = RGBColor(0xC0, 0x00, 0x00)  # 深红标题
ACCENT = RGBColor(0x1A, 0x56, 0xDB)  # 蓝色
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x40)
ORANGE = RGBColor(0xE0, 0x7C, 0x00)
BLUE = RGBColor(0x1A, 0x56, 0xDB)
LBLUE = RGBColor(0xD6, 0xE4, 0xF0)
LRED = RGBColor(0xF8, 0xE0, 0xE0)
LGREEN = RGBColor(0xE0, 0xF0, 0xE4)
LORANGE = RGBColor(0xFD, 0xF0, 0xD5)

def set_bg(slide, color=BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_text(s, l, t, w, h, text, sz=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tf

def add_bullet(tf, text, sz=12, color=BLACK, level=0):
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color; p.level = level

def box(s, l, t, w, h, text, fill=LBLUE, font_color=BLACK, sz=10, bold=False, align=PP_ALIGN.CENTER):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA); shape.line.width = Pt(0.5)
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = font_color; p.font.bold = bold; p.alignment = align
    return shape

def title_bar(s, text, y=0.1):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), Inches(13.33), Inches(0.6))
    shape.fill.solid(); shape.fill.fore_color.rgb = TITLE_BG; shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]; p.text = text; p.font.size = Pt(22); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.LEFT

def arrow_right(s, l, t, w=0.5):
    shape = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(0.25))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

def arrow_down(s, l, t, h=0.4):
    shape = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(0.25), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

# ═══ S1: 封面 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(3))
shape.fill.solid(); shape.fill.fore_color.rgb = TITLE_BG; shape.line.fill.background()
add_text(s, 1, 0.5, 11, 1, "QuantaMind 量智大脑", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(s, 1, 1.6, 11, 0.6, "详细设计方案", 28, WHITE, False, PP_ALIGN.CENTER)
add_text(s, 1, 3.5, 11, 0.5, "量子科技自主科研 AI 中台 · 全生命周期多智能体协同平台", 18, GRAY, False, PP_ALIGN.CENTER)
items = [("12 个 AI Agent", "覆盖全链路"), ("107+ 工具", "Function Calling"), ("13 个开源项目", "四大平台集成"),
         ("3 条流水线", "20/100/105比特"), ("947 条知识", "21份文档索引"), ("4 域 20 表", "跨域追溯")]
for i, (v, d) in enumerate(items):
    x = 0.5 + i * 2.15
    box(s, x, 4.5, 2.0, 1.2, f"{v}\n{d}", LIGHT, BLACK, 12, True)
add_text(s, 1, 6.2, 11, 0.5, "量子科技长三角产业创新中心 · 2026年3月", 14, GRAY, False, PP_ALIGN.CENTER)

# ═══ S2: 战略地位与发展目标 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  战略地位与发展目标")
add_text(s, 0.3, 0.9, 6, 0.4, "战略定位", 18, RED, True)
add_text(s, 0.3, 1.3, 6, 1.5, "QuantaMind 量智大脑是量子芯片全生命周期的 AI 中台底座，\n连接设计→仿真→制造→校准→测控→数据六大环节，\n通过多智能体协同实现自主科研加速。", 13, GRAY)
# 战略框图
box(s, 0.3, 3.0, 2.8, 0.7, "量子芯片设计\nQiskit Metal + KQCircuits", LBLUE, BLACK, 10, True)
arrow_right(s, 3.2, 3.2)
box(s, 3.8, 3.0, 2.2, 0.7, "仿真验证\nHFSS/Q3D/Sonnet", LBLUE, BLACK, 10, True)
arrow_right(s, 6.1, 3.2)
box(s, 6.7, 3.0, 2.2, 0.7, "工艺制造\nCHIPMES + secsgem", LORANGE, BLACK, 10, True)
arrow_right(s, 9.0, 3.2)
box(s, 9.6, 3.0, 2.0, 0.7, "测控校准\nARTIQ + Pulse", LGREEN, BLACK, 10, True)
arrow_right(s, 11.7, 3.2)
box(s, 12.2, 3.0, 0.9, 0.7, "数据\n中台", LRED, BLACK, 9, True)
# 底部：AI 中台
box(s, 0.3, 4.2, 12.8, 0.7, "QuantaMind 量智大脑 · AI 中台（12 Agent · 107 工具 · 知识库 · 流水线 · 自主发现）", LRED, RED, 13, True)
# 右侧目标
add_text(s, 7, 0.9, 6, 0.4, "发展目标", 18, RED, True)
tf = add_text(s, 7, 1.3, 6, 1.5, "· 研发周期缩短 50%+（AI 自动设计+优化迭代）", 12, BLACK)
add_bullet(tf, "· 全链路数据贯通（设计/制造/测控三域实时关联）", 12)
add_bullet(tf, "· 7×24 自主科研（Heartbeat 多级心跳 + 自主发现）", 12)
add_bullet(tf, "· 支持 20/100/105/千比特级芯片设计", 12)
add_bullet(tf, "· 开放生态（12 个开源项目 + 插件化架构）", 12)

# ═══ S3: 系统宏观层级架构图 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  系统宏观层级架构设计")
# 用户层
box(s, 3.5, 0.9, 6, 0.6, "用户交互层：Web 客户端（12页面）· 桌面客户端 · CLI · 项目资料库", LIGHT, BLACK, 11, True)
arrow_down(s, 6.5, 1.55)
# 网关层
box(s, 2, 2.0, 9, 0.6, "网关层 Gateway：FastAPI · 42+ REST API · WebSocket · SSE 流式 · 流水线引擎 · 报告导出", LBLUE, BLUE, 11, True)
arrow_down(s, 6.5, 2.65)
# 智能层 - 三列
box(s, 0.3, 3.1, 3.5, 1.5, "Brain 推理层\n─────────\nSmartBrain(自动降级)\nOpenAICompat(7家LLM)\nOllama(本地)\nFallback(内置)\n─────────\nFunction Calling\n支持 tool_calls", LBLUE, BLACK, 9)
box(s, 4.0, 3.1, 5, 1.5, "Orchestrator 编排层\n─────────────────\n意图路由（关键词匹配 → 12 Agent）\nTool Call 循环（最多 20 轮）\nAgent 专属 system prompt + 工具权限\n对话流水线自动生成\n项目记忆注入", LBLUE, BLACK, 9)
box(s, 9.2, 3.1, 3.8, 1.5, "Memory + Heartbeat\n──────────\nL2 项目记忆(Markdown)\n知识库(947条 · 21文档)\n──────────\nL0 实时(5min)\nL1 巡检(6h)\nL2 实验(12h) · L3 洞察(24h)", LBLUE, BLACK, 9)
arrow_down(s, 6.5, 4.65)
# 工具层
box(s, 0.3, 5.0, 12.7, 0.6, "Hands 工具执行层：107+ 注册工具 · 13 个平台适配器 · 优雅降级（未安装自动 Mock）", LORANGE, BLACK, 11, True)
arrow_down(s, 6.5, 5.65)
# 平台层
box(s, 0.3, 5.9, 3.0, 1.2, "Q-EDA 设计\n────────\nQiskit Metal ✅\nKQCircuits ✅\nQEDA 团队代码\n(16份资料)", LGREEN, BLACK, 9)
box(s, 3.5, 5.9, 3.0, 1.2, "MES 制造\n────────\nCHIPMES (真实)\nOpenMES\nsecsgem · Grafana", LORANGE, BLACK, 9)
box(s, 6.7, 5.9, 3.0, 1.2, "测控\n────────\nARTIQ\nQiskit Pulse\nMitiq(ZNE/PEC/CDR)", LRED, BLACK, 9)
box(s, 9.9, 5.9, 3.1, 1.2, "数据中台\n────────\nQCoDeS · SeaTunnel\nDoris(4域20表)\nqData(治理)", LBLUE, BLACK, 9)

# ═══ S4: 软件全功能及工作流程示意图 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  Q-EDA 软件全功能及工作流程示意图")
# 上半部分：设计流程
steps = [("需求分析", LBLUE), ("架构设计\n拓扑选型", LBLUE), ("参数设计\n频率分配", LBLUE),
         ("仿真计算\nLOM/EPR", LORANGE), ("版图设计\n布线 DRC", LGREEN), ("GDS导出\n掩膜制造", LRED)]
for i, (name, color) in enumerate(steps):
    x = 0.3 + i * 2.15
    box(s, x, 0.9, 1.9, 0.9, name, color, BLACK, 11, True)
    if i < len(steps) - 1:
        arrow_right(s, x + 2.0, 1.2, 0.3)

# 中部：AI 智能体参与
box(s, 0.3, 2.2, 12.7, 0.5, "AI 智能体全程参与：💎芯片设计师 · ⚛️理论物理学家 · 🖥️仿真工程师 · 📚知识工程师 · 📊数据分析师", LRED, RED, 11, True)

# 下半部分：工具调用链
add_text(s, 0.3, 3.0, 12, 0.4, "Tool Call 工作流程（以设计 20 比特芯片为例）", 14, RED, True)
chain = [
    ("1. qeda_catalog\n查阅 QEDA 资料", LBLUE),
    ("2. qeda_get_chip\n获取 20bit 规格", LBLUE),
    ("3. qeda_junction\n获取 JJ 参数", LBLUE),
    ("4. metal_create\n创建 12.5mm 设计", LGREEN),
    ("5. metal_add_transmon\n添加 Q1-Q5 比特", LGREEN),
    ("6. metal_add_route\n路由布线", LGREEN),
    ("7. metal_export_gds\n导出 GDS", LORANGE),
    ("8. 保存到资料库\n+ 数据中台", LRED),
]
for i, (name, color) in enumerate(chain):
    x = 0.15 + (i % 4) * 3.3
    y = 3.5 + (i // 4) * 1.3
    box(s, x, y, 3.0, 1.0, name, color, BLACK, 10)
    if i < len(chain) - 1 and i % 4 < 3:
        arrow_right(s, x + 3.05, y + 0.35, 0.3)

# ═══ S5: 应用层 7 个子系统关系图 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  应用层子系统关系图")
# 上层应用
add_text(s, 0.3, 0.9, 12, 0.3, "上层应用：设计、仿真、制造、测控", 14, RED, True)
box(s, 0.3, 1.3, 3.0, 1.0, "版图设计子系统\nQiskit Metal\nKQCircuits", LBLUE, BLACK, 10, True)
box(s, 3.6, 1.3, 3.0, 1.0, "仿真验证子系统\nHFSS/Q3D/Sonnet\nLOM/EPR 分析", LBLUE, BLACK, 10, True)
box(s, 6.9, 1.3, 3.0, 1.0, "制造协同子系统\nCHIPMES\nOpenMES + secsgem", LORANGE, BLACK, 10, True)
box(s, 10.2, 1.3, 2.8, 1.0, "测控子系统\nARTIQ + Pulse\nMitiq 纠错", LGREEN, BLACK, 10, True)
# 中层核心
add_text(s, 0.3, 2.6, 12, 0.3, "中层核心：数据与智能", 14, RED, True)
box(s, 0.3, 3.0, 6.2, 1.0, "数据管理子系统\nQCoDeS(采集) → SeaTunnel(ETL) → Doris(仓库) → qData(治理)\n4 域 20 表 · 跨域追溯 · AI 训练导出", LORANGE, BLACK, 10, True)
box(s, 6.8, 3.0, 6.2, 1.0, "AI 赋能子系统\nBrain(7 LLM) · Orchestrator(12 Agent 路由)\n107 工具 · 知识库(947条) · Heartbeat(L0-L3)", LRED, BLACK, 10, True)
# 底层支撑
add_text(s, 0.3, 4.3, 12, 0.3, "底层支撑：安全与运维", 14, RED, True)
box(s, 0.3, 4.7, 6.2, 0.8, "权限与安全子系统\nJWT 认证（规划）· RBAC · 审计日志 · 数据分级", LIGHT, GRAY, 10, True)
box(s, 6.8, 4.7, 6.2, 0.8, "系统管理与运维子系统\n日志(ring buffer) · 调试面板 · 配置热更新 · 健康检查", LIGHT, GRAY, 10, True)
# 连线说明
arrow_down(s, 3, 2.35, 0.3); arrow_down(s, 6, 2.35, 0.3); arrow_down(s, 9, 2.35, 0.3); arrow_down(s, 11, 2.35, 0.3)
arrow_down(s, 3, 4.05, 0.3); arrow_down(s, 9, 4.05, 0.3)

# ═══ S6: 发展路线图 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  发展路线图（2026-2028）")
# 时间轴
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.08))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
phases = [
    ("2026 H1\n做优基础", LRED, ["Gateway + Brain + Hands", "12 Agent · 107 工具", "Web 客户端 12 页面", "Q-EDA 集成(Metal+KQC)", "20/105 比特验证"]),
    ("2026 H2\n做优产品", LORANGE, ["CHIPMES 真实对接", "Memory L1(Redis)+L4(RAG)", "安全认证 + 审计", "CI/CD + Docker", "知识图谱 V1.0"]),
    ("2027\n做大市场", LGREEN, ["MoE 路由 + 多模型", "千比特级设计", "全自动流水线", "模型微调", "多团队协同"]),
    ("2028+\n做大生态", LBLUE, ["自主实验设计", "黑灯工厂联动", "开放插件市场", "量子-经典混合", "行业标准制定"]),
]
for i, (title, color, items) in enumerate(phases):
    x = 0.3 + i * 3.25
    box(s, x, 0.9, 3.0, 2.5, "", color, BLACK, 9)
    add_text(s, x + 0.1, 0.95, 2.8, 0.5, title, 14, RED, True, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s, x + 0.15, 1.6 + j * 0.35, 2.7, 0.3, "· " + item, 10, BLACK)
    # 时间轴节点
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.3), Inches(3.65), Inches(0.35), Inches(0.35))
    dot.fill.solid(); dot.fill.fore_color.rgb = RED; dot.line.fill.background()
    # 下方说明
    box(s, x, 4.2, 3.0, 1.5, "", LIGHT, BLACK, 9)
    labels = ["EDA 工具自研\n+ 开源集成", "推动 Q-EDA 生态\n平台化", "做大行业市场\n标准化", "国际化\n开放生态"]
    add_text(s, x + 0.1, 4.3, 2.8, 1.2, labels[i], 12, GRAY, False, PP_ALIGN.CENTER)

# ═══ S7: 12 Agent 协作框图 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  12 AI 科学家 Agent 协作架构")
# 协调者在中间
box(s, 5.2, 1.0, 3.0, 0.8, "🧠 协调者\n意图路由 · 任务分解 · 冲突仲裁", LRED, RED, 10, True)
# 第一排：设计+理论+仿真+算法
agents_row1 = [
    ("💎 芯片设计师", "Metal+KQC·GDS", LBLUE),
    ("⚛️ 理论物理学家", "哈密顿·噪声", LBLUE),
    ("🖥️ 仿真工程师", "HFSS·Q3D", LBLUE),
    ("🔬 算法工程师", "VQE·QAOA", LBLUE),
]
for i, (name, desc, color) in enumerate(agents_row1):
    x = 0.3 + i * 3.25
    box(s, x, 2.2, 3.0, 0.8, f"{name}\n{desc}", color, BLACK, 10, True)
    arrow_down(s, x + 1.3, 1.85, 0.3)
# 第二排：制造+运维+测控+材料
agents_row2 = [
    ("🏭 工艺工程师", "CHIPMES·良率·SPC", LORANGE),
    ("🔧 设备运维员", "ARTIQ·校准·SECS", LORANGE),
    ("📡 测控科学家", "T1/T2·Mitiq纠错", LGREEN),
    ("🧪 材料科学家", "格物·高通量", LGREEN),
]
for i, (name, desc, color) in enumerate(agents_row2):
    x = 0.3 + i * 3.25
    box(s, x, 3.4, 3.0, 0.8, f"{name}\n{desc}", color, BLACK, 10, True)
# 第三排：数据+知识+项目
agents_row3 = [
    ("📊 数据分析师", "Doris·qData·追溯", LORANGE),
    ("📚 知识工程师", "文献·知识图谱·经验", LORANGE),
    ("📋 项目经理", "周报·里程碑·风险", LIGHT),
]
for i, (name, desc, color) in enumerate(agents_row3):
    x = 1.8 + i * 3.25
    box(s, x, 4.6, 3.0, 0.8, f"{name}\n{desc}", color, BLACK, 10, True)
# 底部共享
box(s, 0.3, 5.8, 12.7, 0.6, "共享基础：项目记忆 · 知识库(947条) · 项目资料库(16份QEDA资料) · 工具执行层(107+工具) · 数据中台(4域20表)", LRED, RED, 11, True)

# ═══ S8: 端到端流水线流程图 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  端到端流水线工作流程")
add_text(s, 0.3, 0.85, 12, 0.3, "以 20 比特可调耦合器芯片为例（6 阶段 · 7 Agent · ~101 步 · ~10 秒）", 13, GRAY)
stages = [
    ("阶段1\n芯片设计", "💎", LBLUE, "创建设计\n20 Xmon比特\n19 耦合器\n路由布线\nGDS + 掩膜\n~64 步"),
    ("阶段2\n仿真验证", "🖥️", LBLUE, "Q3D 电容矩阵\nEPR 本征模\nAnsys HFSS\nSonnet 导出\n~8 步"),
    ("阶段3\n工艺制造", "🏭", LORANGE, "8步工艺流程\n批次派工\n良率/SPC\n设备检查\n封装\n~9 步"),
    ("阶段4\n设备校准", "🔧", LGREEN, "ARTIQ 硬件\nPulse 全套校准\n频率+振幅\nDRAG+读出\n~4 步"),
    ("阶段5\n测控表征", "📡", LRED, "光谱/Rabi\nT1/T2/Echo\n读出优化\nMitiq 纠错\n~9 步"),
    ("阶段6\n数据分析", "📊", LORANGE, "跨域追溯\n质量检查\n保存中台\nAI 导出\n~7 步"),
]
for i, (title, emoji, color, desc) in enumerate(stages):
    x = 0.15 + i * 2.17
    box(s, x, 1.3, 2.0, 0.6, title, color, BLACK, 12, True)
    box(s, x, 2.1, 2.0, 3.5, desc, LIGHT, GRAY, 10)
    if i < 5:
        arrow_right(s, x + 2.05, 1.5, 0.2)
# 审批门
add_text(s, 0.3, 5.8, 12, 0.3, "🚦 可选：阶段间人工审批门 | ⏸ 暂停/恢复/终止 | 📄 Word + MD 报告导出 | 每步数据自动保存到数据中台", 12, RED, True)
# 三个模板
box(s, 0.3, 6.3, 4.0, 0.6, "模板1: 20比特 一维链 12.5mm", LBLUE, BLACK, 11, True)
box(s, 4.6, 6.3, 4.0, 0.6, "模板2: 100比特 二维网格 20mm", LORANGE, BLACK, 11, True)
box(s, 8.9, 6.3, 4.1, 0.6, "模板3: 105比特 15×7网格 25mm", LGREEN, BLACK, 11, True)

# ═══ S9: 数据中台框图 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  数据中台与项目资料库架构")
# 数据来源
box(s, 0.3, 1.0, 2.5, 1.0, "Q-EDA\n设计数据\nGDS/参数/仿真", LBLUE, BLACK, 10, True)
box(s, 3.0, 1.0, 2.5, 1.0, "CHIPMES\n制造数据\n工单/良率/SPC", LORANGE, BLACK, 10, True)
box(s, 5.7, 1.0, 2.5, 1.0, "ARTIQ/Pulse\n测控数据\nT1/T2/校准", LGREEN, BLACK, 10, True)
box(s, 8.4, 1.0, 2.5, 1.0, "流水线\n执行数据\n步骤/参数/结果", LRED, BLACK, 10, True)
# 箭头
for x in [1.4, 4.1, 6.8, 9.5]:
    arrow_down(s, x, 2.1, 0.3)
# ETL 层
box(s, 0.3, 2.5, 10.6, 0.6, "SeaTunnel ETL 管道：100+ 连接器 · 批量/实时/CDC 同步 · Exactly-Once", LORANGE, BLACK, 11, True)
arrow_down(s, 5.5, 3.15, 0.3)
# Doris
box(s, 0.3, 3.5, 10.6, 1.5, "", LBLUE, BLACK, 9)
add_text(s, 0.5, 3.55, 10, 0.3, "Apache Doris OLAP 数据仓库 · 亚秒级 SQL · MPP 列式存储", 13, BLUE, True)
box(s, 0.5, 4.0, 2.4, 0.8, "设计域(4表)\ndesigns\nsimulations\ndrc · gds", WHITE, BLACK, 9)
box(s, 3.1, 4.0, 2.4, 0.8, "制造域(5表)\nlots · orders\nyield · spc\nequipment", WHITE, BLACK, 9)
box(s, 5.7, 4.0, 2.4, 0.8, "测控域(5表)\ncharacterization\ncalibration\nbenchmarks", WHITE, BLACK, 9)
box(s, 8.3, 4.0, 2.4, 0.8, "流水线域(6表)\npipeline_runs\nsteps · params\nmeasurements", WHITE, BLACK, 9)
arrow_down(s, 5.5, 5.05, 0.3)
# qData
box(s, 0.3, 5.4, 10.6, 0.6, "qData 数据治理：资产目录 · 标准管理 · 质量检查 · Text2SQL · 血缘追溯 · 数据服务 API", LGREEN, BLACK, 11, True)
# 右侧：资料库
box(s, 11.2, 1.0, 1.9, 5.0, "项目资料库\n────────\n自定义文件夹\n\n📁 QEDA\n  (16份)\n📁 MES数据\n📁 测控记录\n📁 文献\n\n11种格式\n自动解析\n拖拽上传\n持久化索引\n全 Agent\n可调用", LORANGE, BLACK, 9, True)

# ═══ S10: 关键指标 + 总结 ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "  关键技术指标与预期成果")
metrics = [
    ("12", "AI Agent", "覆盖全生命周期"),
    ("107+", "注册工具", "LLM Function Calling"),
    ("42+", "API 端点", "REST + SSE + WS"),
    ("13", "开源项目", "四大平台集成"),
    ("7", "LLM 提供商", "一键热切换"),
    ("947", "知识条目", "21份文档索引"),
    ("3", "流水线模板", "20/100/105比特"),
    ("4域20表", "数据中台", "跨域追溯"),
]
for i, (v, name, desc) in enumerate(metrics):
    x = 0.3 + (i % 4) * 3.25
    y = 1.0 + (i // 4) * 1.6
    box(s, x, y, 3.0, 1.3, f"{v}\n{name}\n{desc}", LIGHT, BLACK, 13, True)
# 预期成果
add_text(s, 0.3, 4.5, 12, 0.4, "预期成果", 18, RED, True)
tf = add_text(s, 0.3, 5.0, 6, 2, "· 填补行业工具空白：量子芯片设计全流程 AI 中台", 13, BLACK)
add_bullet(tf, "· 实现设计与工艺深度融合：设计→制造→测控数据贯通", 13)
add_bullet(tf, "· 提升国产设计工具自主可控能力", 13)
add_bullet(tf, "· 保障核心技术自主可控：降低对 Ansys 等国外软件依赖", 13)
tf2 = add_text(s, 6.5, 5.0, 6, 2, "· 研发效率提升 50%+（AI 自动设计迭代）", 13, BLACK)
add_bullet(tf2, "· 支持千比特级芯片设计（2027 目标）", 13)
add_bullet(tf2, "· 形成开放生态和行业标准", 13)
add_bullet(tf2, "· 持续推进自主科研 AI 系统建设", 13)

out = "E:\\work\\QuantaMind\\docs\\QuantaMind详细设计方案_V2.pptx"
prs.save(out)
print(f"OK: {out}")
