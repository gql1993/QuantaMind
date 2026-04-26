# 项目资料库 — 文件上传、自动解析、数据入库
# 支持格式：Word(.docx) / PDF / Excel(.xlsx) / CSV / GDS/OAS / 图片

import csv
import io
import json
import logging
import os
import re
import shutil
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

from quantamind import config
from quantamind.server import state_store

_log = logging.getLogger("quantamind.project_library")

LIBRARY_DIR = config.DEFAULT_ROOT / "library"
LIBRARY_DIR.mkdir(parents=True, exist_ok=True)

# 文件夹和文件记录（持久化到 JSON）
_folders: List[Dict[str, Any]] = []
_files: List[Dict[str, Any]] = []
_INDEX_FILE = LIBRARY_DIR / "index.json"


def _save_index():
    """将文件夹和文件记录持久化到磁盘"""
    data = {"folders": _folders, "files": [{k: v for k, v in f.items() if k != "parse_result"} for f in _files]}
    try:
        _INDEX_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
    except Exception as e:
        _log.error("保存索引失败: %s", e)


def _load_index():
    """从磁盘恢复文件夹和文件记录"""
    global _folders, _files
    if not _INDEX_FILE.exists():
        return
    try:
        data = json.loads(_INDEX_FILE.read_text(encoding="utf-8"))
        _folders = data.get("folders", [])
        loaded_files = data.get("files", [])
        _files = []
        for f in loaded_files:
            if f.get("path") and os.path.exists(f["path"]):
                f["parse_result"] = None
                _files.append(f)
        _log.info("从索引恢复：%d 个文件夹，%d 个文件", len(_folders), len(_files))
    except Exception as e:
        _log.error("加载索引失败: %s", e)


_load_index()


def _storage_dir(project_id: str, folder_name: str = "") -> Path:
    base_dir = LIBRARY_DIR / project_id
    if folder_name:
        return base_dir / folder_name
    return base_dir


def _get_folder(folder_id: str) -> Optional[Dict[str, Any]]:
    for folder in _folders:
        if folder["folder_id"] == folder_id:
            return folder
    return None


def _normalize_folder_name(name: str) -> str:
    folder_name = str(name or "").strip()
    if not folder_name:
        raise ValueError("文件夹名称不能为空")
    if folder_name in {".", ".."}:
        raise ValueError("文件夹名称不合法")
    if any(sep in folder_name for sep in ("/", "\\", ":", "*", "?", "\"", "<", ">", "|")):
        raise ValueError("文件夹名称不能包含路径或保留字符")
    return folder_name


def _relocate_file_record(record: Dict[str, Any], folder_name: str = "") -> str:
    current_path = Path(str(record.get("path", "")))
    if not current_path.name:
        current_path = _storage_dir(record.get("project_id", "default"), folder_name) / (
            f"{record.get('file_id', 'file')}_{record.get('filename', 'unnamed')}"
        )

    target_dir = _storage_dir(record.get("project_id", "default"), folder_name)
    target_dir.mkdir(parents=True, exist_ok=True)
    target_path = target_dir / current_path.name

    if current_path.exists():
        current_resolved = current_path.resolve()
        target_resolved = target_path.resolve(strict=False)
        if current_resolved != target_resolved:
            shutil.move(str(current_path), str(target_path))
    elif target_path.exists():
        target_path = target_path.resolve(strict=False)
    else:
        raise FileNotFoundError(f"资料库文件不存在: {current_path}")

    record["path"] = str(target_path)
    return record["path"]


def create_folder(name: str, description: str = "", project_id: str = "default") -> Dict[str, Any]:
    try:
        folder_name = _normalize_folder_name(name)
    except ValueError as e:
        return {"error": str(e)}
    if any(f.get("project_id") == project_id and f.get("name") == folder_name for f in _folders):
        return {"error": f"文件夹 {folder_name} 已存在"}

    folder_id = f"D-{uuid.uuid4().hex[:8]}"
    folder = {"folder_id": folder_id, "name": folder_name, "description": description,
              "project_id": project_id, "created_at": _now_iso()}
    _folders.append(folder)
    folder_path = _storage_dir(project_id, folder_name)
    folder_path.mkdir(parents=True, exist_ok=True)
    _save_index()
    return folder


def list_folders(project_id: Optional[str] = None) -> List[Dict]:
    flist = _folders if not project_id else [f for f in _folders if f["project_id"] == project_id]
    result = []
    for f in flist:
        files = [fi for fi in _files if fi.get("folder_id") == f["folder_id"]]
        total_size = sum(fi["size_bytes"] for fi in files)
        result.append({**f, "file_count": len(files), "total_size_bytes": total_size})
    return result


def rename_folder(folder_id: str, name: str) -> Dict:
    folder = _get_folder(folder_id)
    if not folder:
        return {"error": "文件夹不存在"}

    try:
        target_name = _normalize_folder_name(name)
    except ValueError as e:
        return {"error": str(e)}

    if any(
        f["folder_id"] != folder_id
        and f.get("project_id") == folder.get("project_id")
        and f.get("name") == target_name
        for f in _folders
    ):
        return {"error": f"文件夹 {target_name} 已存在"}

    old_name = folder.get("name", "")
    if old_name == target_name:
        return folder

    folder_files = [fi for fi in _files if fi.get("folder_id") == folder_id]
    for fi in folder_files:
        _relocate_file_record(fi, target_name)
        fi["folder_name"] = target_name

    old_dir = _storage_dir(folder.get("project_id", "default"), old_name)
    if old_dir.exists():
        try:
            old_dir.rmdir()
        except OSError:
            pass

    folder["name"] = target_name
    _save_index()
    return {**folder, "moved_files": len(folder_files)}


def delete_folder(folder_id: str) -> Dict:
    for i, f in enumerate(_folders):
        if f["folder_id"] == folder_id:
            folder_dir = _storage_dir(f.get("project_id", "default"), f.get("name", ""))
            # 删除文件夹下所有文件
            to_del = [fi for fi in _files if fi.get("folder_id") == folder_id]
            vectors_removed = 0
            for fi in to_del:
                cleanup = _cleanup_file_artifacts(fi)
                vectors_removed += int(cleanup.get("removed_pgvector", 0) or 0)
                _files.remove(fi)
            if folder_dir.exists():
                try:
                    folder_dir.rmdir()
                except OSError:
                    pass
            _folders.pop(i)
            _save_index()
            return {"deleted": folder_id, "files_deleted": len(to_del), "vectors_removed": vectors_removed}
    return {"error": "文件夹不存在"}


def _now_iso():
    return datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")


def _cleanup_file_artifacts(record: Dict[str, Any]) -> Dict[str, Any]:
    """Delete file artifacts across disk, pgvector and ingest job state."""
    cleanup: Dict[str, Any] = {}
    file_id = str(record.get("file_id", ""))
    filepath = str(record.get("path", ""))

    if filepath and os.path.exists(filepath):
        try:
            os.remove(filepath)
        except Exception as e:
            _log.warning("删除文件失败 %s: %s", filepath, e)

    try:
        from quantamind.server import knowledge_base as kb
        if file_id:
            cleanup = kb.delete_library_record(file_id)
    except Exception as e:
        cleanup = {"error": str(e)}
        _log.warning("删除向量索引失败 %s: %s", file_id, e)

    if file_id:
        job_id = record.get("ingest_job_id") or f"J-del-{file_id}"
        state_store.upsert_library_ingest_job(
            job_id=job_id,
            file_id=file_id,
            filename=record.get("filename", ""),
            project_id=record.get("project_id", "default"),
            status="completed",
            stage="deleted",
            attempts=1,
            payload={"cleanup": cleanup},
        )
    return cleanup


def _run_ingest_pipeline(record: Dict[str, Any], job_id: str, attempts: int = 1) -> Dict[str, Any]:
    state_store.upsert_library_ingest_job(
        job_id=job_id,
        file_id=record["file_id"],
        filename=record["filename"],
        project_id=record["project_id"],
        status="running",
        stage="uploaded",
        attempts=attempts,
        payload={"path": str(record["path"]), "file_type": record["file_type"]},
    )

    try:
        parse_result = parse_file(record)
        record["parsed"] = True
        record["parse_result"] = parse_result
        state_store.upsert_library_ingest_job(
            job_id=job_id,
            file_id=record["file_id"],
            filename=record["filename"],
            project_id=record["project_id"],
            status="running",
            stage="parsed",
            attempts=attempts,
            payload={"parse_type": parse_result.get("type", ""), "file_type": record["file_type"]},
        )
        try:
            from quantamind.server import knowledge_base as kb

            record["vector_index"] = kb.index_library_record(record)
            state_store.upsert_library_ingest_job(
                job_id=job_id,
                file_id=record["file_id"],
                filename=record["filename"],
                project_id=record["project_id"],
                status="completed",
                stage="indexed",
                attempts=attempts,
                payload=record["vector_index"],
            )
        except Exception as e:
            record["vector_index"] = {"error": str(e), "backend": "pgvector"}
            state_store.upsert_library_ingest_job(
                job_id=job_id,
                file_id=record["file_id"],
                filename=record["filename"],
                project_id=record["project_id"],
                status="failed",
                stage="indexed",
                attempts=attempts,
                error_message=str(e),
                payload=record["vector_index"],
            )
    except Exception as e:
        record["parsed"] = False
        record["parse_result"] = {"error": str(e)}
        record["vector_index"] = {"error": str(e)}
        _log.error("解析失败 %s: %s", record["filename"], e)
        state_store.upsert_library_ingest_job(
            job_id=job_id,
            file_id=record["file_id"],
            filename=record["filename"],
            project_id=record["project_id"],
            status="failed",
            stage="parsed",
            attempts=attempts,
            error_message=str(e),
            payload={"file_type": record["file_type"]},
        )

    record["ingest_job_id"] = job_id
    _save_index()
    return record


def save_file(filename: str, content: bytes, project_id: str = "default", folder_id: str = "") -> Dict[str, Any]:
    """保存上传的文件到资料库（可指定文件夹）"""
    file_id = f"F-{uuid.uuid4().hex[:8]}"
    job_id = f"J-{uuid.uuid4().hex[:8]}"
    ext = Path(filename).suffix.lower()
    # 确定存储目录
    folder_name = ""
    if folder_id:
        for f in _folders:
            if f["folder_id"] == folder_id:
                folder_name = f["name"]
                break
    save_dir = _storage_dir(project_id, folder_name)
    save_dir.mkdir(parents=True, exist_ok=True)
    file_path = save_dir / f"{file_id}_{filename}"
    file_path.write_bytes(content)

    record = {
        "file_id": file_id,
        "filename": filename,
        "extension": ext,
        "size_bytes": len(content),
        "project_id": project_id,
        "folder_id": folder_id,
        "folder_name": folder_name,
        "path": str(file_path),
        "uploaded_at": _now_iso(),
        "parsed": False,
        "parse_result": None,
        "file_type": _classify_file(ext),
        "ingest_job_id": job_id,
    }
    _files.append(record)
    _log.info("文件上传: %s (%d bytes) → %s", filename, len(content), file_path)
    return _run_ingest_pipeline(record, job_id=job_id, attempts=1)


def _classify_file(ext: str) -> str:
    type_map = {
        ".docx": "设计文档", ".doc": "设计文档", ".pdf": "图纸/文档",
        ".pptx": "演示文稿", ".ppt": "演示文稿",
        ".xlsx": "参数表格", ".xls": "参数表格", ".csv": "数据文件",
        ".gds": "版图文件", ".oas": "版图文件",
        ".json": "数据文件", ".hdf5": "测量数据", ".h5": "测量数据",
        ".png": "图片", ".jpg": "图片", ".jpeg": "图片",
        ".s2p": "S参数", ".s1p": "S参数",
        ".py": "Python 代码", ".ipynb": "Jupyter 笔记本",
        ".7z": "压缩包", ".zip": "压缩包", ".tar": "压缩包", ".gz": "压缩包",
    }
    return type_map.get(ext, "其他")


def parse_file(record: Dict) -> Dict[str, Any]:
    """根据文件类型自动解析，提取关键数据"""
    ext = record["extension"]
    path = record["path"]

    if ext == ".docx":
        return _parse_docx(path)
    elif ext == ".pdf":
        return _parse_pdf(path)
    elif ext == ".pptx":
        return _parse_pptx(path)
    elif ext == ".xlsx":
        return _parse_excel(path)
    elif ext == ".csv":
        return _parse_csv(path)
    elif ext in (".gds", ".oas"):
        return _parse_gds(path)
    elif ext == ".json":
        return _parse_json(path)
    elif ext == ".py":
        return _parse_python(path)
    elif ext == ".ipynb":
        return _parse_notebook(path)
    elif ext in (".png", ".jpg", ".jpeg"):
        return {"type": "image", "size_bytes": os.path.getsize(path)}
    elif ext in (".7z", ".zip", ".tar", ".gz"):
        return {"type": "archive", "size_bytes": os.path.getsize(path), "note": "压缩包已保存，可下载解压"}
    else:
        return {"type": "unknown", "note": f"文件已保存。{ext} 格式暂不支持自动解析"}


def _parse_docx(path: str) -> Dict:
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        tables_data = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(cells)
            if rows:
                tables_data.append({"headers": rows[0] if rows else [], "rows": rows[1:] if len(rows) > 1 else [], "row_count": len(rows) - 1})

        # 提取数值参数（频率、电感、电容等）
        params = _extract_params_from_text("\n".join(paragraphs))

        return {
            "type": "docx",
            "paragraphs": len(paragraphs),
            "tables": len(tables_data),
            "tables_data": tables_data[:10],
            "text_preview": "\n".join(paragraphs[:20]),
            "extracted_params": params,
        }
    except Exception as e:
        return {"type": "docx", "error": str(e)}


def _parse_pdf(path: str) -> Dict:
    try:
        size = os.path.getsize(path)
        return {"type": "pdf", "size_bytes": size, "note": "PDF 已保存，可在客户端预览"}
    except Exception as e:
        return {"type": "pdf", "error": str(e)}


def _parse_excel(path: str) -> Dict:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            headers = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                row_data = [str(c) if c is not None else "" for c in row]
                if i == 0:
                    headers = row_data
                else:
                    rows.append(row_data)
                if i >= 100:
                    break
            sheets.append({"sheet": name, "headers": headers, "rows": rows[:50], "row_count": len(rows)})
        wb.close()
        return {"type": "excel", "sheets": sheets, "sheet_count": len(sheets)}
    except Exception as e:
        return {"type": "excel", "error": str(e)}


def _parse_csv(path: str) -> Dict:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = list(reader)
        headers = rows[0] if rows else []
        data = rows[1:50] if len(rows) > 1 else []
        # 尝试识别数值列
        numeric_cols = []
        for i, h in enumerate(headers):
            try:
                vals = [float(r[i]) for r in data[:10] if i < len(r) and r[i]]
                if vals:
                    numeric_cols.append({"column": h, "min": min(vals), "max": max(vals), "mean": sum(vals) / len(vals)})
            except (ValueError, IndexError):
                pass
        return {"type": "csv", "headers": headers, "rows": data, "row_count": len(rows) - 1,
                "numeric_summary": numeric_cols}
    except Exception as e:
        return {"type": "csv", "error": str(e)}


def _parse_gds(path: str) -> Dict:
    size = os.path.getsize(path)
    result = {"type": "gds", "size_bytes": size}
    try:
        import klayout.db as pya
        layout = pya.Layout()
        layout.read(path)
        result["cells"] = layout.cells()
        result["top_cells"] = [layout.cell(i).name for i in range(min(layout.cells(), 10))]
        result["layers"] = layout.layer_infos()
    except Exception:
        result["note"] = "GDS 文件已保存。安装 KLayout 后可解析详细层/单元信息"
    return result


def _parse_pptx(path: str) -> Dict:
    """解析 PowerPoint 文件"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        all_text = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts)
            slides.append({"slide": i + 1, "text_preview": slide_text[:200]})
            all_text.append(slide_text)
        params = _extract_params_from_text("\n".join(all_text))
        return {"type": "pptx", "slides": len(prs.slides), "slides_preview": slides[:20],
                "text_preview": "\n---\n".join(all_text[:10])[:2000], "extracted_params": params}
    except Exception as e:
        return {"type": "pptx", "error": str(e), "size_bytes": os.path.getsize(path)}


def _parse_python(path: str) -> Dict:
    """解析 Python 源代码"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            code = f.read()
        lines = code.split("\n")
        # 提取类和函数定义
        classes = [l.strip() for l in lines if l.strip().startswith("class ")]
        functions = [l.strip() for l in lines if l.strip().startswith("def ")]
        imports = [l.strip() for l in lines if l.strip().startswith("import ") or l.strip().startswith("from ")]
        # 提取 docstring
        docstring = ""
        if '"""' in code:
            parts = code.split('"""')
            if len(parts) >= 3:
                docstring = parts[1].strip()[:500]
        return {
            "type": "python", "lines": len(lines), "classes": classes[:10],
            "functions": functions[:20], "imports": imports[:15],
            "docstring": docstring, "code_preview": code[:1500],
        }
    except Exception as e:
        return {"type": "python", "error": str(e)}


def _parse_notebook(path: str) -> Dict:
    """解析 Jupyter Notebook（.ipynb）"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            nb = json.load(f)
        cells = nb.get("cells", [])
        code_cells = [c for c in cells if c.get("cell_type") == "code"]
        md_cells = [c for c in cells if c.get("cell_type") == "markdown"]
        # 提取 markdown 标题
        titles = []
        for c in md_cells:
            src = "".join(c.get("source", []))
            for line in src.split("\n"):
                if line.startswith("#"):
                    titles.append(line.strip())
        # 提取 code 中的 import
        imports = []
        all_code = []
        for c in code_cells:
            src = "".join(c.get("source", []))
            all_code.append(src)
            for line in src.split("\n"):
                l = line.strip()
                if l.startswith("import ") or l.startswith("from "):
                    imports.append(l)
        # 提取参数
        params = _extract_params_from_text("\n".join(all_code))
        kernel = nb.get("metadata", {}).get("kernelspec", {}).get("display_name", "")
        return {
            "type": "notebook", "total_cells": len(cells), "code_cells": len(code_cells),
            "markdown_cells": len(md_cells), "titles": titles[:15], "imports": imports[:15],
            "kernel": kernel, "code_preview": "\n\n".join(all_code[:5])[:2000],
            "extracted_params": params,
        }
    except Exception as e:
        return {"type": "notebook", "error": str(e)}


def _parse_json(path: str) -> Dict:
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict):
            return {"type": "json", "keys": list(data.keys())[:20], "preview": json.dumps(data, ensure_ascii=False)[:500]}
        elif isinstance(data, list):
            return {"type": "json", "items": len(data), "preview": json.dumps(data[:3], ensure_ascii=False)[:500]}
        return {"type": "json"}
    except Exception as e:
        return {"type": "json", "error": str(e)}


def _extract_params_from_text(text: str) -> List[Dict]:
    """从文本中提取量子芯片相关参数"""
    params = []
    patterns = [
        (r"(\d+\.?\d*)\s*GHz", "频率", "GHz"),
        (r"(\d+\.?\d*)\s*MHz", "频率/耦合", "MHz"),
        (r"(\d+\.?\d*)\s*nH", "电感", "nH"),
        (r"(\d+\.?\d*)\s*fF", "电容", "fF"),
        (r"(\d+\.?\d*)\s*[μu]s", "时间", "μs"),
        (r"(\d+\.?\d*)\s*nm", "尺寸", "nm"),
        (r"(\d+\.?\d*)\s*[μu]m", "尺寸", "μm"),
        (r"(\d+\.?\d*)\s*mm", "尺寸", "mm"),
        (r"(\d+\.?\d*)\s*[ΩΩ]", "阻抗", "Ω"),
        (r"(\d+\.?\d*)\s*%", "百分比", "%"),
    ]
    for pattern, category, unit in patterns:
        matches = re.findall(pattern, text)
        for m in matches[:5]:
            params.append({"value": float(m), "unit": unit, "category": category})
    return params[:30]


def list_files(project_id: Optional[str] = None) -> List[Dict]:
    if project_id:
        return _attach_ingest_jobs([f for f in _files if f["project_id"] == project_id])
    return _attach_ingest_jobs(list(_files))


def _file_with_ingest(record: Dict[str, Any]) -> Dict[str, Any]:
    out = dict(record)
    job = state_store.get_library_ingest_job(
        file_id=record.get("file_id"),
        job_id=record.get("ingest_job_id"),
    )
    if job:
        out["latest_ingest_job"] = job
    return out


def _attach_ingest_jobs(records: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not records:
        return []
    jobs = state_store.get_latest_library_ingest_jobs([str(record.get("file_id", "")) for record in records])
    enriched: List[Dict[str, Any]] = []
    for record in records:
        out = dict(record)
        job = jobs.get(str(record.get("file_id", "")))
        if job:
            out["latest_ingest_job"] = job
        enriched.append(out)
    return enriched


def _get_file_record(file_id: str) -> Optional[Dict[str, Any]]:
    for record in _files:
        if record["file_id"] == file_id:
            return record
    return None


def get_file(file_id: str) -> Optional[Dict]:
    record = _get_file_record(file_id)
    if not record:
        return None
    return _file_with_ingest(record)


def move_file(file_id: str, target_folder_id: str) -> Dict:
    """将文件移动到另一个文件夹"""
    f = _get_file_record(file_id)
    if not f:
        return {"error": f"文件 {file_id} 不存在"}
    target_name = ""
    if target_folder_id:
        target_folder = _get_folder(target_folder_id)
        if target_folder:
            target_name = target_folder["name"]
        if not target_name:
            return {"error": f"目标文件夹 {target_folder_id} 不存在"}
    previous_folder = f.get("folder_name") or "未分类"
    _relocate_file_record(f, target_name)
    f["folder_id"] = target_folder_id
    f["folder_name"] = target_name
    _log.info("文件移动: %s %s → %s", f["filename"], previous_folder, target_name or "未分类")
    _save_index()
    return {
        "moved": file_id,
        "filename": f["filename"],
        "target_folder": target_name or "未分类",
        "path": f["path"],
    }


def retry_file(file_id: str) -> Dict[str, Any]:
    record = _get_file_record(file_id)
    if not record:
        return {"error": f"文件 {file_id} 不存在"}
    if not record.get("path") or not os.path.exists(record["path"]):
        return {"error": "源文件不存在，无法重试"}

    latest_job = state_store.get_library_ingest_job(
        file_id=record.get("file_id"),
        job_id=record.get("ingest_job_id"),
    )
    attempts = int((latest_job or {}).get("attempts", 0) or 0) + 1
    record["parsed"] = False
    record["parse_result"] = None
    record.pop("latest_ingest_job", None)
    job_id = f"J-{uuid.uuid4().hex[:8]}"
    _log.info("文件重试解析/索引: %s (%s), attempt=%d", record["filename"], file_id, attempts)
    updated = _run_ingest_pipeline(record, job_id=job_id, attempts=attempts)
    return _file_with_ingest(updated)


def retry_folder(folder_id: str, failed_only: bool = True) -> Dict[str, Any]:
    folder = _get_folder(folder_id)
    if not folder:
        return {"error": "文件夹不存在"}

    files = _attach_ingest_jobs([f for f in _files if f.get("folder_id") == folder_id])
    if failed_only:
        targets = [f for f in files if (f.get("latest_ingest_job") or {}).get("status") == "failed"]
    else:
        targets = files

    results = []
    for file in targets:
        results.append(retry_file(file["file_id"]))

    success_count = sum(1 for item in results if not item.get("error"))
    failed_count = len(results) - success_count
    return {
        "folder_id": folder_id,
        "folder_name": folder.get("name", ""),
        "requested": len(targets),
        "retried": success_count,
        "failed": failed_count,
        "results": results[:20],
    }


def delete_file(file_id: str) -> Dict:
    """删除资料库中的文件"""
    for i, f in enumerate(_files):
        if f["file_id"] == file_id:
            cleanup = _cleanup_file_artifacts(f)
            _files.pop(i)
            _log.info("文件删除: %s (%s)", f["filename"], file_id)
            _save_index()
            return {"deleted": file_id, "filename": f["filename"], "cleanup": cleanup}
    return {"error": f"文件 {file_id} 不存在"}


def delete_all_files(project_id: str = "default") -> Dict:
    """清空项目下的所有文件"""
    to_remove = [f for f in _files if f["project_id"] == project_id]
    vectors_removed = 0
    for f in to_remove:
        cleanup = _cleanup_file_artifacts(f)
        vectors_removed += int(cleanup.get("removed_pgvector", 0) or 0)
        _files.remove(f)
    _save_index()
    return {"deleted_count": len(to_remove), "project_id": project_id, "vectors_removed": vectors_removed}


def search_files(query: str) -> List[Dict]:
    q = query.lower()
    matches = [
        f for f in _files if q in f["filename"].lower() or q in f.get("file_type", "").lower()
        or q in json.dumps(f.get("parse_result", {}), ensure_ascii=False).lower()
    ]
    return _attach_ingest_jobs(matches)


def get_stats() -> Dict:
    total_size = sum(f["size_bytes"] for f in _files)
    by_type = {}
    for f in _files:
        t = f["file_type"]
        by_type[t] = by_type.get(t, 0) + 1
    return {"total_files": len(_files), "total_size_bytes": total_size, "by_type": by_type,
            "parsed": sum(1 for f in _files if f["parsed"])}
